#!/usr/bin/env python3
"""Extract text from engagement source documents with PII redaction.

Routes input documents into two destinations:

* **Engagement docs** (CDC, mode opératoire) — committed in-repo under
  ``docs/wip/intelligence/<engagement-slug>/extracts/`` after redaction.
* **Inspiration sources** (competitor distillation, transcripts not yet
  routed) — written to an off-repo working directory under
  ``~/.cache/holistika/inspiration/`` so the raw competitor methodology can
  be distilled before deletion. Source identifiers are operator-supplied via
  ``--redact-token`` and never enter the committed extract.

Supported formats:

* ``.pdf``  via pdfplumber
* ``.docx`` via python-docx (paragraphs + tables)
* ``.pptx`` via python-pptx (slides + shape text)
* ``.xlsx`` via openpyxl (cells + formulas, values + types)
* ``.json`` passthrough (already-extracted workbook dump)
* ``.md``, ``.txt`` passthrough (already-redacted-by-source transcripts)

PII redaction is **deliberately conservative**: a configurable token-list
driven from ``--redact-token`` (repeatable) plus standard email and phone
masks. Operators add explicit redaction tokens; the script never tries to
infer names from arbitrary text.

Cross-references:
* ``akos-governance-remediation.mdc`` — SoT/DI/SOC/DX
* ``akos-mirror-template.mdc`` — redaction discipline
* an earlier one-off operator workbook-serialisation helper; this script
  generalises that pattern into a governed engagement-extraction tool.

Usage::

    py scripts/extract_engagement_pdfs.py engagement \
        --input "<path-to-CDC.docx>" \
        --slug "2026-05-10-suez-webuy-procure-to-pay" \
        --redact-token "[BRIDGE-COLLABORATOR-NAME]"

    py scripts/extract_engagement_pdfs.py inspiration \
        --input "<path-to-competitor-cost-model.json>" \
        --redact-token "[COMPETITOR-NAME]" --redact-token "[COMPETITOR-INTERNAL-CODE]"

Exit codes::

    0 — extraction completed for every input file.
    1 — one or more inputs failed to parse.
"""

from __future__ import annotations

import argparse
import json
import logging
import os
import re
import sys
import unicodedata
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Sequence

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from akos.io import REPO_ROOT
from akos.log import setup_logging

logger = logging.getLogger("akos.engagement_extract")

EMAIL_REGEX = re.compile(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}")
PHONE_REGEX = re.compile(
    r"(?<!\w)(?:\+?\d{1,3}[\s.-]?)?(?:\(\d{1,4}\)[\s.-]?)?\d{2,4}(?:[\s.-]?\d{2,4}){2,4}(?!\w)"
)


@dataclass
class ExtractionConfig:
    """Operator-supplied controls for an extraction run."""

    mode: str  # "engagement" | "inspiration"
    inputs: list[Path]
    slug: str | None
    redact_tokens: list[str] = field(default_factory=list)
    out_root_override: Path | None = None

    def out_root(self) -> Path:
        if self.out_root_override is not None:
            return self.out_root_override
        if self.mode == "engagement":
            if not self.slug:
                raise ValueError("engagement mode requires --slug")
            return (
                REPO_ROOT
                / "docs"
                / "wip"
                / "intelligence"
                / self.slug
                / "extracts"
            )
        cache_env = os.environ.get("HOLISTIKA_INSPIRATION_CACHE", "").strip()
        if cache_env:
            return Path(cache_env).expanduser()
        return Path.home() / ".cache" / "holistika" / "inspiration"


def redact_text(text: str, tokens: Sequence[str]) -> str:
    """Apply email + phone + token redaction."""
    if not text:
        return text
    redacted = text
    for token in tokens:
        if not token:
            continue
        pattern = re.compile(re.escape(token), re.IGNORECASE)
        redacted = pattern.sub("[REDACTED]", redacted)
    redacted = EMAIL_REGEX.sub("[REDACTED-EMAIL]", redacted)
    redacted = PHONE_REGEX.sub("[REDACTED-PHONE]", redacted)
    return redacted


def _format_block(title: str, body: str) -> str:
    return f"\n\n=== {title} ===\n\n{body.rstrip()}\n"


def extract_pdf(path: Path) -> str:
    import pdfplumber

    chunks: list[str] = []
    with pdfplumber.open(path) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            chunks.append(_format_block(f"page {idx:03d}", text))
            tables = page.extract_tables() or []
            for t_idx, table in enumerate(tables, start=1):
                rendered = "\n".join(
                    " | ".join((cell or "").strip() for cell in row) for row in table
                )
                chunks.append(_format_block(f"page {idx:03d} table {t_idx}", rendered))
    return "".join(chunks)


def extract_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    out: list[str] = []
    for p_idx, para in enumerate(document.paragraphs, start=1):
        if para.text.strip():
            style = para.style.name if para.style and para.style.name else "Normal"
            out.append(f"[{style}] {para.text.strip()}")
    if document.tables:
        for t_idx, table in enumerate(document.tables, start=1):
            out.append(f"\n[Table {t_idx}]")
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                out.append(" | ".join(cells))
    return "\n".join(out)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[str] = []
    for s_idx, slide in enumerate(prs.slides, start=1):
        body: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            body.append(run.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    body.append(" | ".join(cells))
        chunks.append(_format_block(f"slide {s_idx:03d}", "\n".join(body)))
    return "".join(chunks)


def extract_xlsx(path: Path) -> str:
    import openpyxl
    from openpyxl.utils import get_column_letter

    wb_values = openpyxl.load_workbook(path, data_only=True)
    wb_formulas = openpyxl.load_workbook(path, data_only=False)
    chunks: list[str] = []
    for sheet_name in wb_values.sheetnames:
        ws_v = wb_values[sheet_name]
        ws_f = wb_formulas[sheet_name]
        rows: list[str] = []
        max_row = ws_v.max_row or 0
        max_col = ws_v.max_column or 0
        for r in range(1, max_row + 1):
            cells: list[str] = []
            any_filled = False
            for c in range(1, max_col + 1):
                v = ws_v.cell(row=r, column=c).value
                f = ws_f.cell(row=r, column=c).value
                addr = f"{get_column_letter(c)}{r}"
                if v is None and f is None:
                    cells.append("")
                    continue
                any_filled = True
                if isinstance(f, str) and f.startswith("="):
                    cells.append(f"{addr}={f}~={v!r}")
                else:
                    cells.append(f"{addr}={v!r}")
            if any_filled:
                rows.append(" | ".join(cells))
        chunks.append(_format_block(f"sheet {sheet_name}", "\n".join(rows)))
    return "".join(chunks)


def extract_json(path: Path) -> str:
    """Render JSON to a flat key=value text representation."""
    payload = json.loads(path.read_text(encoding="utf-8"))
    return _format_block(f"json {path.name}", json.dumps(payload, ensure_ascii=False, indent=2))


def extract_passthrough(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def extract_one(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".xlsx":
        return extract_xlsx(path)
    if suffix == ".json":
        return extract_json(path)
    if suffix in {".md", ".txt"}:
        return extract_passthrough(path)
    raise ValueError(f"unsupported extension: {suffix} (path={path})")


def safe_stem(path: Path) -> str:
    """Filesystem-safe stem (lowercase, ascii-ish)."""
    stem = path.stem
    stem = stem.replace(" ", "_").replace("-", "-").replace("é", "e").replace("è", "e").replace("à", "a")
    stem = re.sub(r"[^A-Za-z0-9._-]+", "_", stem)
    return stem.lower()


def write_extract(out_dir: Path, source: Path, body: str) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    target = out_dir / f"{safe_stem(source)}.txt"
    header = (
        f"# extract: {source.name}\n"
        f"# bytes_in: {source.stat().st_size}\n"
        f"# chars_out: {len(body)}\n"
        f"# redaction: emails+phones+operator-tokens\n\n"
    )
    target.write_text(header + body, encoding="utf-8")
    return target


def resolve_path(path: Path) -> Path:
    """Return *path* if it exists, else try NFD/NFC normalisations.

    Windows + macOS hand back filenames with combining accents (NFD) while
    operator-supplied paths often arrive pre-composed (NFC). The extractor
    accepts whichever form actually resolves on disk.
    """
    if path.exists():
        return path
    parent = path.parent
    if not parent.exists():
        return path
    name = path.name
    candidates = {
        unicodedata.normalize("NFC", name),
        unicodedata.normalize("NFD", name),
        unicodedata.normalize("NFKC", name),
        unicodedata.normalize("NFKD", name),
    }
    try:
        listing = os.listdir(parent)
    except OSError:
        return path
    for entry in listing:
        if entry in candidates or unicodedata.normalize("NFC", entry) == unicodedata.normalize("NFC", name):
            return parent / entry
    return path


def run_extract(config: ExtractionConfig) -> int:
    failures = 0
    out_root = config.out_root()
    out_root.mkdir(parents=True, exist_ok=True)
    logger.info("Mode: %s | output: %s", config.mode, out_root)
    for raw_path in config.inputs:
        path = resolve_path(raw_path)
        if not path.exists():
            logger.error("missing: %s", raw_path)
            failures += 1
            continue
        try:
            raw = extract_one(path)
            redacted = redact_text(raw, config.redact_tokens)
            written = write_extract(out_root, path, redacted)
            logger.info("ok: %s -> %s (%d chars)", path.name, written, len(redacted))
        except Exception as exc:
            logger.error("failed: %s (%s)", path, exc)
            failures += 1
    return failures


def parse_args(argv: list[str] | None = None) -> ExtractionConfig:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("mode", choices=("engagement", "inspiration"))
    parser.add_argument("--input", action="append", required=True, help="input file (repeatable)")
    parser.add_argument("--slug", help="engagement slug (engagement mode only)")
    parser.add_argument(
        "--redact-token",
        action="append",
        default=[],
        help="token to redact (case-insensitive; repeatable)",
    )
    parser.add_argument("--out-root", help="override output directory (testing)")
    parser.add_argument("--json-log", action="store_true")
    args = parser.parse_args(argv)
    setup_logging(json_output=args.json_log)
    return ExtractionConfig(
        mode=args.mode,
        inputs=[Path(i) for i in args.input],
        slug=args.slug,
        redact_tokens=list(args.redact_token),
        out_root_override=Path(args.out_root) if args.out_root else None,
    )


def main(argv: list[str] | None = None) -> int:
    config = parse_args(argv)
    return run_extract(config)


if __name__ == "__main__":
    sys.exit(main())
